import concurrent.futures
import datetime
import io
import os
import shutil
from typing import List

import PyPDF2
from bson import ObjectId
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sumy.nlp.tokenizers import Tokenizer
from sumy.parsers.plaintext import PlaintextParser
from sumy.summarizers.text_rank import TextRankSummarizer

from app.config import Config
from app.models.mongoClient import MongoClient
from app.utils import Common, Subscription, files_and_folders
from app.utils.email_helper import send_mail
from app.utils.formatter import cursor_to_dict, get_base64_encoding
from app.utils.llm.llm_highlights import generate_highlights
from app.utils.pipelines import PipelineStages
from app.utils.production import Production
from app.utils.socket import (emit_document_upload_status, socket_error,
                              socket_info, socket_success)
from app.utils.vectorstore.document_loaders import DocumentLoader

from . import user_service as UserService


class MyDocumentsService:
    
    @staticmethod
    def upload_document(logged_in_user, file, path: str, upload_id: str):
        """
        The `upload_document` function uploads a file to a specified path, parses and inserts the
        document into a database, updates the virtual filename, and saves the file on disk or a cloud
        bucket.

        Args:
          logged_in_user: The logged_in_user parameter is an object that represents the currently logged
        in user. It contains information about the user, such as their ID, name, email, etc.
          file: The `file` parameter is the file object that is being uploaded. It contains the actual
        file data and metadata such as the filename.
          path: The `path` parameter represents the directory path where the document should be
        uploaded. It is a string that specifies the location within the file system or cloud storage
        where the document should be saved.

        Returns:
          a tuple containing two values: 1) an integer indicating the success or failure of the upload
        process (1 for success, 0 for failure), and 2) the inserted ID of the document in the database.
        """
        user_id = str(logged_in_user["_id"])
        print("File upload path : ", path)
        new_path = "/" + logged_in_user["_id"] + "/"
        if path != "/":
            new_path += path[1:]
        filename = file.filename
        print(f"File {filename}")
        file_extension = filename.rsplit(".", 1)[-1]

        if file_extension not in ["pdf", "doc", "docx", "pptx", "ppt", "txt"]:
            socket_info(
                user_id,
                f"Skipping upload of {filename} due to incompatible file format"
            )
            emit_document_upload_status(user_id, upload_id, f"Skipping upload of {filename} due to incompatible file format", 20)
            return 0, None

        # Parse and insert document into database
        inserted_id = MyDocumentsService().parse_document(
            logged_in_user, file, new_path
        )

        if not inserted_id:
            socket_error(
                user_id,
                f"Failed to save {filename} to database due to some error..."
            )
            emit_document_upload_status(user_id, upload_id, f"Failed to save {filename} to database due to some error...", 20)
            return 0, None

        # Update the virtual filename of the file based on its inserted ID
        update_response = MyDocumentsService().update_virtual_filename(
            inserted_id, file_extension
        )

        # Save file on disk or cloud bucket
        MyDocumentsService().save_file(file, inserted_id, logged_in_user, path)
        
        # After a document has been successfully uploaded change document subscription
        print("🧾 Updating subscription: Document used...")
        file_size = files_and_folders.get_size([file])
        UserService.update_document_subscription(logged_in_user["_id"], file_size)

        return 1, inserted_id

    @staticmethod
    def upload_documents(logged_in_user, files, path: str, upload_id: str):
        """
        The function `upload_documents` uploads multiple files to a specified path, using a
        ThreadPoolExecutor to execute the upload process concurrently, and updates the document
        vectorstore for each successfully uploaded file.

        Args:
          logged_in_user: The logged_in_user parameter is the user object of the currently logged in
        user. It contains information about the user, such as their ID, name, email, etc.
          files: The `files` parameter is a list of files that you want to upload. Each file should be
        in a format that can be processed by the `upload_document` method of the `MyDocumentsService`
        class.
          path: The `path` parameter is the directory path where the documents will be uploaded to.
        """
        user_id = str(logged_in_user["_id"])
        
        emit_document_upload_status(user_id, upload_id, f"Uploading {len(files)} document(s)...", 20)

        # Create a ThreadPoolExecutor with a specified number of threads (e.g., 4)
        with concurrent.futures.ThreadPoolExecutor(max_workers=8) as executor:
            # Submit the function with arguments to the thread pool
            results = [
                executor.submit(
                    MyDocumentsService().upload_document, logged_in_user, file, path, upload_id, 
                )
                for file in files
            ]

        # Getting function returns from all function calls from threadpool
        emit_document_upload_status(user_id, upload_id, f"Saved documents to database...", 50)
        outputs = [result.result() for result in results]
        print("MyDocumentsService outputs : ", outputs)
        # All document _ids inserted
        uploaded_documents_ids = [output[1] for output in outputs if output[1]]
        print("MyDocumentsService uploaded_documents_ids : ", uploaded_documents_ids)
        # Number of documents successfully uploaded
        uploaded_documents_num = sum([output[0] for output in outputs])
        print("MyDocumentsService uploaded_documents_num : ", uploaded_documents_num)

        # Update document vectorstore for each successfully inserted id
        emit_document_upload_status(user_id, upload_id, f"Parsing documents and getting them ready for chat...", 70)
        for inserted_id in uploaded_documents_ids:
            file = MyDocumentsService().get_file(inserted_id)
            # print("MyDocumentsService file : ", file)
            virtual_file_name = file["virtualFileName"]
            filepath = MyDocumentsService().get_file_save_path(
                virtual_file_name, user_id, path
            )
            print("MyDocumentsService filepath : ", filepath)
            DocumentLoader(user_id, inserted_id, file, filepath).load_document()

        # Calculating number of documents successfully uploaded
        if uploaded_documents_num > 0:
            socket_success(
                user_id, f"Successfully uploaded {uploaded_documents_num} documents!"
            )
            emit_document_upload_status(user_id, upload_id, f"Successfully uploaded {uploaded_documents_num} documents!", 100)

    @staticmethod
    def get_file_save_path(filename, user, path):
        """
        The function `get_file_save_path` returns the save path for a file based on the filename, user,
        and path provided.

        Args:
          filename: The name of the file that you want to save.
          user: The `user` parameter is the user ID or identifier. It is used to determine if the file
        is created by the user or if it is shared with the user.
          path: The `path` parameter is the path where the file should be saved. It is a string
        representing the directory structure where the file should be stored.

        Returns:
          the file save path.
        """
        file = MyDocumentsService().get_file_by_virtual_name(filename)
        file_created_by = file["createdBy"]["_id"]
        if str(user) == str(file_created_by):
            user_folder_path = os.path.join(Config.USER_FOLDER, str(user))
            new_path = path[1:]
            if path != None:
                user_folder_path = os.path.join(user_folder_path, new_path)
            file_save_path = os.path.join(user_folder_path, filename)
        else:
            replace_substring = "/" + str(file_created_by) + "/"
            user_folder_path = os.path.join(Config.USER_FOLDER, str(file_created_by))
            file_root = file["root"]
            file_root = file_root.replace(replace_substring, "")
            user_folder_path = os.path.join(user_folder_path, file_root)
            file_save_path = os.path.join(user_folder_path, filename)
        return file_save_path

    @staticmethod
    def get_file_path(file, user_id):
        """
        The function `get_file_path` takes a file object and a user ID as input, and returns the file
        path based on the root folder and virtual file name.

        Args:
          file: The "file" parameter is a dictionary that contains information about a file. It has two
        keys: "root" and "virtualFileName".
          user_id: The user ID is a unique identifier for each user. It is used to identify the user's
        folder where the file is located.

        Returns:
          the file path as a string if it is successfully generated. If there is an exception, it will
        return None.
        """
        try:
            root = file["root"]
            root_folder_path = os.path.join(Config.USER_FOLDER, user_id)
            folder_substring = "/" + user_id + "/"
            file_name = file["virtualFileName"]
            if root == folder_substring:
                file_path = os.path.join(root_folder_path, file_name)
            else:
                new_root = root.replace(folder_substring, "")
                folder = os.path.join(root_folder_path, new_root)
                file_path = os.path.join(folder, file_name)

            # print("File path : ", file_path)

            return file_path

        except Exception as e:
            Common.exception_details("my_documents_service.get_file_path", e)
            return None

    def parse_document(self, logged_in_user, file, new_path):
        """
        The function `parse_document` takes in a logged-in user, a file, and a new path, and based on
        the file extension, it calls different parsing functions to process the file and returns an
        inserted ID.

        Args:
          logged_in_user: The logged_in_user parameter is the user who is currently logged in and
        performing the document parsing operation.
          file: The `file` parameter is the file object that represents the document to be parsed. It is
        passed to the `parse_document` method as an argument.
          new_path: The `new_path` parameter is the path where the parsed document will be saved. It
        specifies the location where the parsed document will be stored after it has been processed.

        Returns:
          the variable "inserted_id".
        """
        try:
            filename = file.filename
            file_extension = filename.rsplit(".", 1)[-1]

            # PDF
            if file_extension == "pdf":
                print("Parsing pdf...")
                inserted_id = self._parse_pdf(file, filename, logged_in_user, new_path)

            # WORD
            elif file_extension == "doc" or file_extension == "docx":
                print("Parsing doc/docx...")
                inserted_id = self._parse_doc(file, filename, logged_in_user, new_path)

            # PPT
            elif file_extension == "pptx":
                print("Parsing pptx...")
                inserted_id = self._parse_pptx(file, filename, logged_in_user, new_path)

            # TXT
            elif file_extension == "txt":
                print("Parsing txt...")
                inserted_id = self._parse_text(file, filename, logged_in_user, new_path)

            else:
                print("Failed to parse invalid file format...")
                inserted_id = None

            return inserted_id

        except Exception as e:
            Common.exception_details("my_documents_service.parse_document", e)
            return None

    def update_virtual_filename(self, file_id, file_extension):
        """
        Updates the virtual filename of the given file to "file_id.file_extension"
        """
        m_db = MongoClient.connect()

        virtual_filename = file_id + "." + file_extension
        response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].update_one(
            {"_id": ObjectId(file_id)}, {"$set": {"virtualFileName": virtual_filename}}
        )
        return response.modified_count

    def get_all_files(self, user_id, root, limit: int = 20, offset: int = 0):
        """
        Retrieves all the files uploaded by the user
        """
        m_db = MongoClient.connect()
        common_pipeline = MyDocumentsService._get_my_documents_pipeline()
        new_root = "/" + user_id + "/"

        if root and root != "/":
            new_root += root[1:]

        uploaded_documents = [
            PipelineStages.stage_match(
                {"createdBy._id": ObjectId(user_id), "root": new_root}
            )
        ] + common_pipeline

        shared_documents = [
            PipelineStages.stage_match(
                {"usersWithAccess": {"$in": [ObjectId(user_id)]}}
            ),
            PipelineStages.stage_lookup(
                Config.MONGO_USER_MASTER_COLLECTION,
                "createdBy._id",
                "_id",
                "userDetails",
            ),
            PipelineStages.stage_unwind("userDetails"),
            PipelineStages.stage_add_fields({"owner": "$userDetails.name"}),
            PipelineStages.stage_unset(["userDetails", "usersWithAccess"]),
        ] + common_pipeline

        # Use the facet pipeline stage to get both uploaded and shared documents
        pipeline = [
            PipelineStages.stage_facet(
                {"uploaded": uploaded_documents, "shared": shared_documents}
            ),
            {"$skip": offset},
            {"$limit": limit},
        ]

        documents = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].aggregate(pipeline)
        return cursor_to_dict(documents)

    def get_file(self, file_id):
        """
        The function `get_file` retrieves a file from a MongoDB database based on its ID.

        Args:
          file_id: The `file_id` parameter is the unique identifier of the file that you want to
        retrieve from the database.

        Returns:
          the first document that matches the given file_id.
        """
        m_db = MongoClient.connect()

        pipeline = [
            PipelineStages.stage_match({"_id": ObjectId(file_id)})
        ] + MyDocumentsService._get_my_documents_pipeline()

        response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].aggregate(pipeline)

        return cursor_to_dict(response)[0]
    
    def get_files_by_id(self, file_ids):
        """
        The function `get_file` retrieves files from a MongoDB database based on their IDs.

        Args:
        file_ids: The `file_ids` parameter is a list of unique identifiers of the files that you want to
        retrieve from the database.

        Returns:
        A list of documents that match the given file_ids.
        """
        m_db = MongoClient.connect()

        pipeline = [
            PipelineStages.stage_match({"_id": {"$in": [ObjectId(file_id) for file_id in file_ids]}})
        ] + MyDocumentsService._get_my_documents_pipeline()

        response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].aggregate(pipeline)
        return cursor_to_dict(response)

    def rename_document(self, _id, rename_value, user_id):
        """
        The `rename_document` function renames a document (file or folder) in a database, and if the
        document is a folder, it also renames the corresponding folder in a cloud storage service (GCP)
        or in the system.

        Args:
          _id: The `_id` parameter is the unique identifier of the document that needs to be renamed. It
        is used to query the document from the database.
          rename_value: The `rename_value` parameter is the new name that you want to assign to the
        document.
          user_id: The `user_id` parameter is the unique identifier of the user who created the
        document.

        Returns:
          the number of modified documents in the database.
        """
        m_db = MongoClient.connect()
        query = {
            "_id": ObjectId(_id),
            "createdBy._id": ObjectId(user_id),
        }
        doc = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].find_one(query)
        # See if the document is a folder
        if doc["type"] == "Folder":
            root = doc["root"]

            # If we are on GCP then rename the folder in GCP
            if Config.GCP_PROD_ENV:
                if root == "/" + str(user_id) + "/":
                    folder_to_rename = root[1:] + doc["originalFileName"] + "/"
                    new_folder_name = root[1:] + rename_value + "/"
                else:
                    folder_to_rename = root[1:] + "/" + doc["originalFileName"] + "/"
                    new_folder_name = root[1:] + "/" + rename_value + "/"
                # print("Folder to rename : ", folder_to_rename)
                # print("New folder name: ", new_folder_name)

                bucket = Production.get_users_bucket()
                source_blob = bucket.blob(folder_to_rename)
                destination_blob = bucket.blob(new_folder_name)
                # Creating destination folder
                destination_blob.upload_from_string("")
                # print("Source blob : ", source_blob)
                # print("Destination blob : ", destination_blob)

                # list all blobs in source folder
                blobs = bucket.list_blobs(prefix=folder_to_rename)
                # print("Listing all blobs in source folder : ")
                for blob in blobs:
                    blob_name = str(blob.name)
                    new_blob_name = blob_name.replace(folder_to_rename, new_folder_name)
                    destination_blob = bucket.blob(new_blob_name)
                    # Copy the file to the new new_blob_name
                    blob_copy = bucket.copy_blob(blob, bucket, new_blob_name)
                    # Delete the source blob after copying
                    blob.delete()
                old_root_value = "/" + folder_to_rename[:-1]
                new_root_value = "/" + new_folder_name[:-1]
            # If we are on not on GCP then rename the folder in the system
            else:
                username_substring = "/" + user_id + "/"
                new_root = root.replace(username_substring, "")
                if root != username_substring:
                    new_root = new_root + "/"
                folder_to_rename = new_root + doc["originalFileName"]
                new_folder_name = new_root + rename_value
                user_folder_path = os.path.join(Config.USER_FOLDER, user_id)
                old_folder_path = os.path.join(user_folder_path, folder_to_rename)
                new_folder_path = os.path.join(user_folder_path, new_folder_name)

                try:
                    os.rename(old_folder_path, new_folder_path)
                except Exception as e:
                    print(f"Error renaming {folder_to_rename}")
                    print("================================================")
                    print(e)
                    print("================================================")
                    return None
                old_root_value = "/" + folder_to_rename
                new_root_value = "/" + new_folder_name
                # print("Old root value: " + old_root_value)
                # print("New root value: " + new_root_value)

            # Update the root of all the files and folders in this current folder
            children_records = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].find(
                {"root": {"$regex": old_root_value}}
            )
            res_count = 0
            for child in children_records:
                child_root = child["root"]
                updated_child_root = child_root.replace(old_root_value, new_root_value)
                update_child = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].update_one(
                    {
                        "_id": child["_id"],
                    },
                    {"$set": {"root": updated_child_root}},
                )
                res_count += update_child.modified_count

            # Update the original file name field for that record in the database
            response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].update_one(
                query,
                {"$set": {"originalFileName": rename_value}},
            )
            return response.modified_count + res_count
        # Since the document is a file, we need to update the database with the <new-name>.extension
        else:
            virtualFileName = doc["virtualFileName"]
            extension = virtualFileName[virtualFileName.rfind(".") + 1 :]
            # Update the original file name field for that record in the database
            response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].update_one(
                query,
                {"$set": {"originalFileName": rename_value + "." + extension}},
            )
            return response.modified_count

    def get_all_files_by_virtual_name(self, user_id, virtual_file_names):
        """
        The function retrieves all files by virtual name for a specific user from a MongoDB collection.
        
        :param user_id: The `user_id` parameter is the unique identifier of the user for whom we want to
        retrieve files. It is used to filter the files based on the user who created them
        :param virtual_file_names: The `virtual_file_names` parameter is a list of virtual file names
        that you want to retrieve from the database for a specific user identified by `user_id`. The
        function `get_all_files_by_virtual_name` uses these parameters to query the database and return
        the documents that match the criteria specified in the
        :return: The function `get_all_files_by_virtual_name` is returning the result of the MongoDB
        aggregation query performed on the collection specified in
        `Config.MONGO_DOCUMENT_MASTER_COLLECTION`. The result is being converted from a cursor to a
        dictionary before being returned.
        """
        m_db = MongoClient.connect()

        pipeline = [
            PipelineStages.stage_match({
                "createdBy._id": ObjectId(user_id),
                "virtualFileName": {"$in": virtual_file_names},
            })
        ] + MyDocumentsService._get_my_documents_pipeline()

        response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].aggregate(pipeline)

        return cursor_to_dict(response)

    def get_file_by_virtual_name(self, virtual_name):
        """
        The function retrieves a document from a MongoDB collection based on its virtual file name.

        Args:
          virtual_name: The virtual name is a parameter that represents the name of the file you want to
        retrieve from the database.

        Returns:
          the document that matches the given virtual name from the specified collection in the MongoDB
        database.
        """
        m_db = MongoClient.connect()

        document = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].find_one(
            {"virtualFileName": virtual_name}
        )

        return document

    def save_file(self, original_file, file_id, user, path):
        """
        The `save_file` function saves a file either to a cloud storage bucket or to a local folder,
        depending on the value of the `Config.GCP_PROD_ENV` variable.

        Args:
          original_file: The original file that needs to be saved.
          file_id: The `file_id` parameter is the unique identifier of the file that needs to be saved.
        It is used to retrieve information about the file from the database.
          user: The "user" parameter is a dictionary that contains information about the user. It likely
        includes details such as the user's ID, name, email, and other relevant information.
          path: The `path` parameter is a string that represents the path where the file should be
        saved. It can be an absolute path or a relative path.
        """
        key = "_id"
        user_id = user[key]
        # Get virtual filename from DB
        file = self.get_file(file_id)
        virtual_file_name = file["virtualFileName"]
        file_root = file["root"]

        if Config.GCP_PROD_ENV:
            # print("PATH : ", path)
            # print("Original file :", original_file)

            if path == "/":
                folder_name = str(user_id) + path
            else:
                folder_name = file_root[1:] + "/"

            bucket = Production.get_users_bucket()
            file_blob = bucket.blob(folder_name + virtual_file_name)

            # finding the file extension
            extension = virtual_file_name[virtual_file_name.rfind(".") + 1 :]

            original_file.stream.seek(0)
            if extension == "pdf":
                file_blob.upload_from_string(
                    original_file.read(), content_type="application/pdf"
                )
            else:
                file_blob.upload_from_string(original_file.read())

            print(f"File {virtual_file_name} uploaded to {folder_name} successfully.")
            print("Path sent to get_file_save_path: ", path)

        else:
            # Ensure that the user image upload folder exists
            folder_path = os.path.join(Config.USER_FOLDER, str(user_id))
            os.makedirs(folder_path, exist_ok=True)

            # Save file
            file_save_path = self.get_file_save_path(virtual_file_name, user_id, path)
            original_file.stream.seek(0)
            original_file.save(file_save_path)

            print("Saved file!")

    # FOLDER CODE
    def create_folder(self, folder_name, path, user_id):
        """
        The `create_folder` function creates a folder in a specified path for a user and saves the
        folder data in a database.

        Args:
          folder_name: The name of the folder to be created.
          path: The `path` parameter represents the directory path where the folder should be created.
        It can be an absolute path or a relative path. If it is "/", it means the folder should be
        created in the root directory.
          user_id: The user_id parameter is the unique identifier of the user who is creating the
        folder.

        Returns:
          the inserted ID of the folder in the database if the folder is successfully created. If the
        folder is not created, it returns None.
        """
        # print("PATH : " + path)
        folder_created_flag = False
        m_db = MongoClient.connect()
        db_path = ""
        if Config.GCP_PROD_ENV:
            if path == "/":
                folder_to_create = str(user_id) + path + folder_name + "/"
                db_path = "/" + user_id + "/"
            else:
                folder_to_create = str(user_id) + path + "/" + folder_name + "/"
                db_path = "/" + user_id + path
            # print("CREATE FOLDER : " + folder_to_create)

            # Create an empty blob to represent the folder (blobs are like objects in GCS)
            bucket = Production.get_users_bucket()
            folder_blob = bucket.blob(folder_to_create)
            # Upload an empty string to create an empty folder
            folder_blob.upload_from_string("")
            folder_created_flag = True
            print(f"Created folder {folder_name} in bucket!")
        else:
            if path == "/":
                user_folder_path = os.path.join(Config.USER_FOLDER, user_id)
                print(f"Creating folder {folder_name} in {user_folder_path}!")
                db_path = "/" + user_id + "/"
            else:
                new_path = path[1:]
                # print("NEW PATH : ", new_path)
                user_folder_path = os.path.join(Config.USER_FOLDER, user_id)
                user_folder_path = os.path.join(user_folder_path, new_path)
                print(f"Creating folder {folder_name} in {user_folder_path}!")
                db_path = "/" + user_id + path

            if folder_name != None:
                folder_save_path = os.path.join(user_folder_path, folder_name)
            else:
                folder_save_path = user_folder_path
            # print("Folder save path : ", folder_save_path)
            if not os.path.exists(folder_save_path):
                # print("Folder to make : ", folder_save_path)
                os.makedirs(folder_save_path)
                folder_created_flag = True

        # If folder is created :
        if folder_created_flag == True:
            folder_data = {
                "originalFileName": folder_name,
                "createdBy": {"_id": ObjectId(user_id), "ref": "user"},
                "createdOn": datetime.datetime.utcnow(),
                "type": "Folder",
                "root": db_path,
                # "root": "/" + ObjectId(user_id) + "/" if path == None else path
            }
            # print("DB PATH : " + db_path)
            response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].insert_one(
                folder_data
            )

            if response:
                return str(response.inserted_id)
        else:
            return None

    def delete_folder(self, folder_id, user_id):
        """
        The `delete_folder` function deletes a folder and its contents from a system database and file
        storage.

        Args:
          folder_id: The ID of the folder to be deleted.
          user_id: The user_id parameter is the unique identifier of the user who owns the folder that
        needs to be deleted.

        Returns:
          the value of the variable "folder_delete_flag", which indicates whether the folder was
        successfully deleted or not.
        """

        m_db = MongoClient.connect()

        # Deleting folder from system database
        folder = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].find_one(
            {"_id": ObjectId(folder_id), "type": "Folder"}
        )
        db_path = "/" + user_id + "/"
        folder_name = folder["originalFileName"]
        root = folder["root"]
        root_folder_path = os.path.join(Config.USER_FOLDER, user_id)
        # print("ROOT :", root)
        if root == db_path:
            folder_path = os.path.join(root_folder_path, folder_name)
        else:
            new_root = root.replace(db_path, "")
            folder_path = os.path.join(root_folder_path, new_root)
            folder_path = os.path.join(folder_path, folder_name)
        # print("FOLDER PATH :", folder_path)
        folder_delete_flag = False
        # Delete folder contents
        if Config.GCP_PROD_ENV:
            try:
                # print(root)
                if root != db_path:
                    folder_to_del = root[1:] + "/" + folder_name
                    # print("Folder : " + folder_to_del)
                else:
                    folder_to_del = root[1:] + folder_name
                    # print("folder_to_delete is %s" % folder_to_del)
                bucket = Production.get_users_bucket()
                blobs = bucket.list_blobs(prefix=folder_to_del + "/")
                for blob in blobs:
                    blob.delete()
                print(f"Folder '{folder_name}' deleted from bucket successfully.")
                folder_delete_flag = True
            except Exception as e:
                print(f"Error deleting folder with name '{folder_name}': {e}")
                folder_delete_flag = False
        else:
            if os.path.isdir(folder_path):
                # Iterate over all the files and subdirectories in the folder
                for filename in os.listdir(folder_path):
                    file_path = os.path.join(folder_path, filename)
                    # Check if it's a file or a directory
                    if os.path.isfile(file_path):
                        # Delete the file
                        os.remove(file_path)
                    else:
                        # Recursively delete the subdirectory
                        shutil.rmtree(file_path)

                # Delete empty folder
                os.rmdir(folder_path)
                folder_delete_flag = True
                print(f"Deleted folder: {folder_path}")
            else:
                print(f"Not a valid folder: {folder_path}")
                folder_delete_flag = False

        if folder_delete_flag == True:
            # fetch root using folder_id
            # rooexp=concat folder name with root
            # ex: '/userid/foldername'
            # delete all where root like rooexp

            # print("Folder :", folder_name)
            # print("Root : ", root)
            if root == "/" + user_id + "/":
                folder_to_del = root + folder_name
            else:
                folder_to_del = root + "/" + folder_name
            print("Deleting ... " + folder_to_del)
            query = {
                "root": {"$regex": folder_to_del},
                "createdBy._id": ObjectId(user_id),
            }
            matching_documents = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].find(
                query
            )
            # for doc in matching_documents:
            #     # print(doc)
            #     # print()
            result = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].delete_many(query)
            del_main_folder = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].delete_one(
                {
                    "_id": ObjectId(folder_id),
                    "createdBy._id": ObjectId(user_id),
                    "type": "Folder",
                }
            )
            print("Deleted folder from database")
        else:
            print(f"Couldn't delete folder!")

        return folder_delete_flag
        # Delete document if the user is the one who uploaded the file
        # delete_response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].delete_one({
        #     "_id": ObjectId(folder_id),
        #     "createdBy._id": ObjectId(user_id)
        # })

        # if result.deleted_count > 0:
        #     return True
        # else:
        #     return False

    def get_all_folders(self, user_id):
        """
        The function `get_all_folders` retrieves all folders created by a specific user from a MongoDB
        database and returns them as a list of dictionaries.

        Args:
          user_id: The user_id parameter is the unique identifier of the user for whom we want to
        retrieve all the folders.

        Returns:
          a list of folders that match the given user_id. Each folder in the list is a dictionary with
        modified values. The "_id" field and "createdBy._id" field are converted to strings. If the
        folder has a "usersWithAccess" field, each value in the list is converted to a string.
        """
        m_db = MongoClient.connect()
        query = {"createdBy._id": ObjectId(user_id), "type": "Folder"}

        folders = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].find(query)
        result = []
        for folder in list(folders):
            folder["_id"] = str(folder["_id"])
            folder["createdBy"]["_id"] = str(folder["createdBy"]["_id"])
            if "usersWithAccess" in folder.keys():
                length = len(folder["usersWithAccess"])
                if length > 0:
                    for j in range(length):
                        folder["usersWithAccess"][j] = str(folder["usersWithAccess"][j])
            result.append(folder)

        return result

    def move_file_to_folder(self, file, folder, user_id):
        """
        The `move_file_to_folder` function moves a file to a specified folder location.

        Args:
          file: The `file` parameter is a dictionary that represents the file to be moved. It contains
        information such as the virtual file name, root path, and other details.
          folder: The "folder" parameter represents the destination folder where the file will be moved
        to. It can be a string representing the folder path or an object containing information about
        the folder, such as its root and original file name.
          user_id: The user ID is a unique identifier for a user in the system. It is used to locate the
        user's files and folders in the file system or database.

        Returns:
          a boolean value indicating whether the file was successfully moved or not.
        """
        m_db = MongoClient.connect()

        file_moved_flag = False
        file_save_root = ""
        substring = "/" + user_id + "/"
        virtualFileName = file["virtualFileName"]

        # Get the folder's root
        if folder == "/":
            folder_root = substring
        else:
            folder_root = folder["root"]

        if Config.GCP_PROD_ENV:
            file_root = file["root"]
            if file_root == substring:
                source_blob_name = file_root[1:] + virtualFileName
            else:
                source_blob_name = file_root[1:] + "/" + virtualFileName
            # print("File root : ", file_root)
            if folder == "/":
                file_save_root = folder_root
                destination_blob_name = folder_root[1:] + virtualFileName
            else:
                # print("Folder root : ", folder_root)
                file_save_root = folder_root + folder["originalFileName"]
                destination_blob_name = (
                    folder_root[1:] + folder["originalFileName"] + "/" + virtualFileName
                )
            # print("Destination blob name : ", destination_blob_name)
            # source_blob_name = "https://storage.cloud.google.com/" + self.bucket_name + "/" +  source_blob_name
            # print("Source blob name : ", source_blob_name)

            bucket = Production.get_users_bucket()
            source_blob = bucket.blob(source_blob_name)
            destination_blob = bucket.blob(destination_blob_name)
            # print("Source blob : ", source_blob)
            # print("Destination blob : ", destination_blob)
            # finding the file extension
            # extension = virtualFileName[virtualFileName.rfind('.')+1:]
            # if(extension == 'pdf'):
            #     destination_blob.upload_from_filename(source_blob_name, content_type="application/pdf")
            # else:
            #     destination_blob.upload_from_filename(source_blob_name)

            # Copy the file to the new location
            destination_generation_match_precondition = 0
            # source_blob.stream.seek(0)
            blob_copy = bucket.copy_blob(
                source_blob,
                bucket,
                destination_blob_name,
                if_generation_match=destination_generation_match_precondition,
            )
            file_moved_flag = True
            # print("BLOB COPY : ", blob_copy)
            # source_blob.copy_to(destination_blob)

            # Copy the source blob to the destination
            # token = None
            # while True:
            #     token, bytes_rewritten, total_bytes = source_blob.rewrite(destination_blob, token=token)
            #     print("Bytes written : ", bytes_rewritten)
            #     print("Total bytes written : ", total_bytes)
            #     if token is None:
            #         file_moved_flag = False
            #         break
            #     else:
            #         file_moved_flag = True

            # Delete the source blob after copying
            source_blob.delete()

            # print(
            #     f"File '{source_blob_name}' moved to '{destination_blob_name}' in bucket '{self.bucket_name}'."
            # )

        else:
            root_folder_path = os.path.join(Config.USER_FOLDER, user_id)

            # Locating file in original path using it's root value
            file_root = file["root"]
            if file_root == ("/" + user_id + "/"):
                file_path = os.path.join(root_folder_path, virtualFileName)
            else:
                new_file_root = file_root.replace(substring, "")
                file_folder_path = os.path.join(root_folder_path, new_file_root)
                file_path = os.path.join(file_folder_path, virtualFileName)
            # Moving the file to the new location
            if folder_root == substring:
                if folder == "/":
                    destination_folder = root_folder_path
                else:
                    destination_folder = os.path.join(
                        root_folder_path, folder["originalFileName"]
                    )
                if os.path.isfile(file_path):
                    shutil.move(file_path, destination_folder)
                    file_moved_flag = True
                    file_save_root = folder_root
                    if folder != "/":
                        file_save_root += folder["originalFileName"]

                    print(
                        f"Moved file {virtualFileName} from '{file_path}' to '{destination_folder}'"
                    )
                else:
                    file_moved_flag = False
                    print(f"Could not move file {virtualFileName} as its not a file")
            else:
                new_folder_root = folder_root.replace(substring, "")
                folder_path = os.path.join(root_folder_path, new_folder_root)
                destination_folder = os.path.join(
                    folder_path, folder["originalFileName"]
                )
                print("DESTINATION PATH: %s" % destination_folder)
                if os.path.isfile(file_path):
                    shutil.move(file_path, destination_folder)
                    file_moved_flag = True
                    file_save_root = folder_root
                    if folder != "/":
                        file_save_root += "/" + folder["originalFileName"]

                    print(
                        f"Moved file {virtualFileName} from '{file_path}' to '{destination_folder}'"
                    )
                else:
                    file_moved_flag = False
                    print(f"Could not move file {virtualFileName} as its not a file")

        if file_moved_flag:
            # See if there is another file with the same name in the new folder
            print("File save root : ", file_save_root)
            files_in_root = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].find(
                {"root": file_save_root}
            )
            # for file in files_in_root:
            #     print("File : ", file)
            # Create a set of existing filenames for efficient lookup
            existing_filenames = set(
                curr_file["originalFileName"] for curr_file in files_in_root
            )

            # If the filename already exists, increment the index count
            index = 1
            new_filename = file["originalFileName"]
            while new_filename in existing_filenames:
                # new_filename = f"{filename}({index})"
                # last_index = filename.rfind(".")
                file_name_without_extension, file_extension = os.path.splitext(
                    file["originalFileName"]
                )
                # print("File name without extension:", file_name_without_extension)
                new_filename = f"{file_name_without_extension}({index}){file_extension}"
                index = index + 1

            update_file_root_response = m_db[
                Config.MONGO_DOCUMENT_MASTER_COLLECTION
            ].update_one(
                {"_id": ObjectId(file["_id"])},
                {"$set": {"root": file_save_root, "originalFileName": new_filename}},
            )
            return bool(update_file_root_response.modified_count)
        else:
            return False

    def get_folder_contents(self, user_id, _id):
        """
        The function `get_folder_contents` retrieves the contents of a folder based on the provided user
        ID and folder ID.

        Args:
          user_id: The user ID is a unique identifier for a user in the system. It is used to identify
        the user who is requesting the folder contents.
          _id: The `_id` parameter is the unique identifier of the folder in the database. It is used to
        retrieve the folder document from the database.

        Returns:
          the contents of a folder in the form of a dictionary.
        """

        m_db = MongoClient.connect()
        common_pipeline = MyDocumentsService._get_my_documents_pipeline()

        folder = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].find_one(
            {"_id": ObjectId(_id)}
        )
        print("FOLDER : ", folder)
        root = folder["root"]
        substring = "/" + str(folder["createdBy"]["_id"]) + "/"
        if root == substring:
            new_root = root + folder["originalFileName"]
        else:
            new_root = root + "/" + folder["originalFileName"]
        pipeline = [PipelineStages.stage_match({"root": new_root})] + common_pipeline

        documents = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].aggregate(pipeline)

        return cursor_to_dict(documents)

    def my_documents_summarize_gpt(self, document_ids: list, sentence_count: int):
        """
        Summarizes the descriptions of user-uploaded documents one-by-one using Text Rank method from sumy module
        Sends them back one by one
        """
        m_db = MongoClient.connect()

        # Convert kiIds to a list of ObjectIDs
        document_ids = list(map(ObjectId, document_ids))

        documents = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].find(
            {"_id": {"$in": document_ids}}
        )

        for idx, document in enumerate(documents):
            document_description = ""
            document_title = document["title"]
            document_created_date = str(document["createdOn"])

            if not document:
                yield document_description, idx, document_title, document_created_date, False

            itemized_summary = document["itemizedSummary"]

            yield itemized_summary, idx, document_title, document_created_date, True

    def highlights(self, file_ids: List[str]):
        """
        The `highlights` function takes a list of file IDs, retrieves the corresponding files from a
        MongoDB collection, generates a summary for each file if it doesn't already exist, and yields
        the summary, title, and index for each file.

        Args:
          file_ids (List[str]): The `file_ids` parameter is a list of strings that represent the IDs of
        the files you want to retrieve highlights for.
        """
        files_collection = Config.MONGO_DOCUMENT_MASTER_COLLECTION
        print(f"fileids: {file_ids}")
        m_db = MongoClient.connect()

        # Convert kiIds to a list of ObjectIDs
        file_ids_object_id = list(map(ObjectId, file_ids))

        files_iterable = m_db[files_collection].find(
            {"_id": {"$in": file_ids_object_id}}
        )

        for idx, file in enumerate(files_iterable):
            description = ""
            title = file["title"]
            # print("Title : ", title)
            if not file:
                print("File not found!")
                yield generate_highlights(description), title, idx

            description = file["description"]
            # If the summary of the document was already generated and stored then return that summary
            highlights_summary = file["highlightsSummary"]
            if len(highlights_summary) > 0:
                print("Found summary from DB!")
                text_summary = highlights_summary

            else:
                print("Generating summary!")
                text_summary = generate_highlights(description)
                # Store the text_summary for next time quick access
                m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].update_one(
                    {"_id": ObjectId(file["_id"])},
                    {"$set": {"highlightsSummary": text_summary}},
                )
            # yield text_summary, True
            yield text_summary, title, idx

    def create_itemized_summary_ppt(self, request_body):
        """
        The function `create_itemized_summary_ppt` creates a PowerPoint presentation with itemized
        summaries based on the given request body.

        Args:
          request_body: The `request_body` parameter is a list of dictionaries. Each dictionary
        represents an item to be included in the itemized summary PowerPoint presentation. Each
        dictionary should have two keys: "title" and "summary". The "title" key should contain the title
        of the item, and the "summary"

        Returns:
          the file path of the created PowerPoint presentation.
        """
        # import Presentation class
        # from pptx library

        """ Ref for slide types:
            0 -> title and subtitle
            1 -> title and content
            2 -> section header
            3 -> two content
            4 -> Comparison
            5 -> Title only
            6 -> Blank
            7 -> Content with caption
            8 -> Pic with caption
            """

        # Creating presentation object
        root = Presentation()

        # Set slide width and height
        slide_width = 11887200
        slide_height = 6686550
        root.slide_width = slide_width
        root.slide_height = slide_height

        # Heading Page
        #############################

        # Creating slide layout
        first_slide_layout = root.slide_layouts[5]

        # Creating slide object to add
        slide = root.slides.add_slide(first_slide_layout)

        # Set slide dimensions
        slide.slide_width = slide_width
        slide.slide_height = slide_height

        left = top = 0
        img_path = "app/static/ppt_images/cover_export.jpg"
        pic = slide.shapes.add_picture(
            img_path, left, top, width=slide_width, height=slide_height
        )

        # This moves it to the background
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # Adding title and subtitle in slide
        # i.e. first page of slide
        title_shape = slide.shapes.title
        title_shape.text = "Itemized Summary"
        # Set title font to Roboto
        title_shape.text_frame.paragraphs[0].font.name = "Roboto"
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(
            255, 255, 255
        )  # Set title color to white
        title_shape.text_frame.paragraphs[0].font.size = Pt(43)  # Set title font size

        # Change the position of the title box
        title_shape.left = Inches(0.85)
        title_shape.top = Inches(2.52)
        title_shape.width = Inches(4)  # Set the width to 4 inches
        title_shape.height = Inches(2)  # Set the width to 2 inches
        title_shape.text_frame.paragraphs[0].font.bold = True

        # Content pages
        ##########################

        # Function to add a new slide with given title and content

        def add_new_slide(title, content):
            slide_layout = root.slide_layouts[5]
            slide = root.slides.add_slide(slide_layout)
            slide.slide_width = slide_width
            slide.slide_height = slide_height

            # Add background image
            img_path = "app/static/ppt_images/inside_export.jpg"
            pic = slide.shapes.add_picture(
                img_path, 0, 0, width=slide_width, height=slide_height
            )

            # This moves it to the background
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)

            # Set title text
            title_shape = slide.shapes.title
            title_shape.text = title.strip()
            # title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # Set title color to white
            title_shape.text_frame.paragraphs[0].font.size = Pt(
                24
            )  # Set title font size
            # Set title font to Roboto
            title_shape.text_frame.paragraphs[0].font.name = "Roboto"
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            title_shape.text_frame.paragraphs[0].font.bold = True

            # Move 0.25 inches to the right horizontally
            title_shape.left = Inches(0.28)
            title_shape.top = Inches(0.26)  # Move 0.16 inches down vertically
            title_shape.width = Inches(8)  # Set the width to 4 inches
            title_shape.height = Inches(0.5)  # Set the width to 0.75 inches

            textbox = slide.shapes.add_textbox(
                Inches(0.28), Inches(1.0), Inches(12.44), Inches(5.3)
            )
            textbox.text_frame.paragraphs[0].text = content.strip()
            textbox.text_frame.paragraphs[0].font.name = "Roboto"
            textbox.text_frame.paragraphs[0].font.size = Pt(13)
            textbox.text_frame.paragraphs[0].font.bold = True
            textbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(72, 88, 95)
            textbox.text_frame.word_wrap = True

            textbox2 = slide.shapes.add_textbox(
                Inches(10.4), Inches(6.7), Inches(1.5), Inches(0.5)
            )
            textbox2.text_frame.paragraphs[0].text = "Itemized Summary"
            textbox2.text_frame.paragraphs[0].font.name = "Roboto"
            textbox2.text_frame.paragraphs[0].font.size = Pt(11)
            textbox2.text_frame.paragraphs[0].font.color.rgb = RGBColor(72, 88, 95)
            textbox2.text_frame.word_wrap = True

        # Check if the example text fits on a single slide
        max_chars_per_slide = 3000
        for file in request_body:
            print(file)
            title = file["title"]
            summary = file["summary"]
            if len(summary) <= max_chars_per_slide:
                # Text fits on a single slide
                summary = summary.strip()
                add_new_slide(title, summary)
            else:
                # Text does not fit on a single slide, split it into multiple slides
                sentences = summary.split(".")
                current_slide_text = ""
                for sentence in sentences:
                    if (
                        len(current_slide_text) + len(sentence) + 1
                        > max_chars_per_slide
                    ):
                        # Current slide is full, add a new slide with the accumulated text
                        current_slide_text = current_slide_text.strip()
                        add_new_slide(title, current_slide_text)
                        current_slide_text = ""
                    current_slide_text += sentence + ". "
                if current_slide_text:
                    # Add the remaining text to a new slide
                    current_slide_text = current_slide_text.strip()
                    add_new_slide(title, current_slide_text)

        # Thank You page
        ######################
        slide_layout = root.slide_layouts[5]  # Title and content layout
        slide = root.slides.add_slide(slide_layout)
        slide.slide_width = slide_width
        slide.slide_height = slide_height

        # Add background image
        img_path = "app/static/ppt_images/thank_you.jpg"
        pic = slide.shapes.add_picture(
            img_path, 0, 0, width=slide_width, height=slide_height
        )

        # This moves it to the background
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # Set title text
        title_shape = slide.shapes.title
        title_shape.text = "Thank you for the download"
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(
            89, 100, 116
        )  # Set title color to white
        title_shape.text_frame.paragraphs[0].font.size = Pt(22)  # Set title font size
        title_shape.text_frame.paragraphs[0].font.bold = True
        # Set title font to Roboto
        title_shape.text_frame.paragraphs[0].font.name = "Roboto"

        # Move 0.2 inches to the right horizontally
        title_shape.left = Inches(4.24)
        title_shape.top = Inches(4.19)  # Move 3 inches down vertically
        title_shape.width = Inches(4.52)  # Set the width to 4 inches
        title_shape.height = Inches(0.45)  # Set the width to 2 inches
        # Saving file

        timestamp = datetime.datetime.utcnow().timestamp()
        # root.save(f"app/ppt_outputs/Output_{timestamp}.pptx")
        # return f"ppt_outputs\Output_{timestamp}.pptx"
        file_name = f"itemized_summary_{timestamp}.pptx"

        # Ensure that the user image upload folder exists
        os.makedirs(Config.USER_SUMMARY_PPTX_DOWNLOAD_FOLDER, exist_ok=True)

        file_path = os.path.join(Config.USER_SUMMARY_PPTX_DOWNLOAD_FOLDER, file_name)
        root.save(file_path)

        return file_path

    def create_highlights_summary_ppt(self, request_body):
        """
        Creates a ppt of the highlights summary and returns its path
        """

        # import Presentation class
        # from pptx library

        """ Ref for slide types:
            0 -> title and subtitle
            1 -> title and content
            2 -> section header
            3 -> two content
            4 -> Comparison
            5 -> Title only
            6 -> Blank
            7 -> Content with caption
            8 -> Pic with caption
            """
        # Creating presentation object
        root = Presentation()

        # Set slide width and height
        slide_width = 11887200
        slide_height = 6686550
        root.slide_width = slide_width
        root.slide_height = slide_height

        # Heading Page
        #############################

        # Creating slide layout
        first_slide_layout = root.slide_layouts[5]

        # Creating slide object to add
        slide = root.slides.add_slide(first_slide_layout)

        # Set slide dimensions
        slide.slide_width = slide_width
        slide.slide_height = slide_height

        left = top = 0
        img_path = "app/static/ppt_images/cover_export.jpg"
        pic = slide.shapes.add_picture(
            img_path, left, top, width=slide_width, height=slide_height
        )

        # This moves it to the background
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # Adding title and subtitle in slide
        # i.e. first page of slide
        title_shape = slide.shapes.title
        title_shape.text = "Highlights Summary"
        # Set title font to Roboto
        title_shape.text_frame.paragraphs[0].font.name = "Roboto"
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(
            255, 255, 255
        )  # Set title color to white
        title_shape.text_frame.paragraphs[0].font.size = Pt(43)  # Set title font size

        # Change the position of the title box
        title_shape.left = Inches(0.85)
        title_shape.top = Inches(2.52)
        title_shape.width = Inches(4)  # Set the width to 4 inches
        title_shape.height = Inches(2)  # Set the width to 2 inches
        title_shape.text_frame.paragraphs[0].font.bold = True

        # Content pages
        ##########################

        # Function to add a new slide with given title and content

        def add_new_slide(title, content):
            slide_layout = root.slide_layouts[5]
            slide = root.slides.add_slide(slide_layout)
            slide.slide_width = slide_width
            slide.slide_height = slide_height

            # Add background image
            img_path = "app/static/ppt_images/inside_export.jpg"
            pic = slide.shapes.add_picture(
                img_path, 0, 0, width=slide_width, height=slide_height
            )

            # This moves it to the background
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)

            # Set title text
            title_shape = slide.shapes.title
            title_shape.text = title.strip()
            # title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # Set title color to white
            title_shape.text_frame.paragraphs[0].font.size = Pt(
                24
            )  # Set title font size
            # Set title font to Roboto
            title_shape.text_frame.paragraphs[0].font.name = "Roboto"
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            title_shape.text_frame.paragraphs[0].font.bold = True

            # Move 0.25 inches to the right horizontally
            title_shape.left = Inches(0.28)
            title_shape.top = Inches(0.26)  # Move 0.16 inches down vertically
            title_shape.width = Inches(8)  # Set the width to 4 inches
            title_shape.height = Inches(0.5)  # Set the width to 0.75 inches

            textbox = slide.shapes.add_textbox(
                Inches(10.4), Inches(6.7), Inches(2), Inches(0.5)
            )
            textbox.text_frame.paragraphs[0].text = "Highlights Summary"
            textbox.text_frame.paragraphs[0].font.name = "Roboto"
            textbox.text_frame.paragraphs[0].font.size = Pt(11)
            textbox.text_frame.paragraphs[0].font.bold = True
            textbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(72, 88, 95)
            textbox.text_frame.word_wrap = True

            textbox2 = slide.shapes.add_textbox(
                Inches(0.28), Inches(1.26), Inches(12), Inches(5.3)
            )
            text_frame = textbox2.text_frame
            for value in enumerate(content):
                keys = list(value[1].keys())
                text = value[1][keys[1]]
                keywords = value[1][keys[0]]

                bullet_point = "\u2022"

                allText = bullet_point + " " + text
                if keywords != "":
                    sentence = allText.split(keywords.lower())
                # else:
                #     sentence = allText.split(keywords)
                para = text_frame.add_paragraph()
                iIndex = 0
                for phrase in sentence:
                    run1 = para.add_run()
                    run1.text = phrase
                    run1.font.name = "Roboto"
                    run1.font.size = Pt(13)

                    if iIndex != len(sentence) - 1:
                        run2 = para.add_run()
                        run2.text = keywords.lower()
                        run2.font.name = "Roboto"
                        run2.font.size = Pt(13)
                        run2.font.color.rgb = RGBColor(0, 0, 255)

                    iIndex = iIndex + 1
                #  Add bullet points to the text box
            textbox2.text_frame.word_wrap = True

        def divide_array(arr, subarray_length):
            divided_array = []
            for i in range(0, len(arr), subarray_length):
                divided_array.append(arr[i : i + subarray_length])
            return divided_array

        for file in request_body:
            title = file["title"]
            highlights = file["highlights"]
            highlights_sub_array = divide_array(highlights, 15)
            for i in highlights_sub_array:
                add_new_slide(title, i)

        # Thank You page
        ######################
        slide_layout = root.slide_layouts[5]  # Title and content layout
        slide = root.slides.add_slide(slide_layout)
        slide.slide_width = slide_width
        slide.slide_height = slide_height

        # Add background image
        img_path = "app/static/ppt_images/thank_you.jpg"
        pic = slide.shapes.add_picture(
            img_path, 0, 0, width=slide_width, height=slide_height
        )

        # This moves it to the background
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # Set title text
        title_shape = slide.shapes.title
        title_shape.text = "Thank you for the download"
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(
            89, 100, 116
        )  # Set title color to white
        title_shape.text_frame.paragraphs[0].font.size = Pt(22)  # Set title font size
        title_shape.text_frame.paragraphs[0].font.bold = True
        # Set title font to Roboto
        title_shape.text_frame.paragraphs[0].font.name = "Roboto"

        # Move 0.2 inches to the right horizontally
        title_shape.left = Inches(4.24)
        title_shape.top = Inches(4.19)  # Move 3 inches down vertically
        title_shape.width = Inches(4.52)  # Set the width to 4 inches
        title_shape.height = Inches(0.45)  # Set the width to 2 inches
        # Saving file

        timestamp = datetime.datetime.utcnow().timestamp()
        # root.save(f"app/ppt_outputs/Output_{timestamp}.pptx")
        # return f"ppt_outputs\Output_{timestamp}.pptx"
        file_name = f"Itemized_summary_{timestamp}.pptx"

        # Ensure that the user image upload folder exists
        os.makedirs(Config.USER_SUMMARY_PPTX_DOWNLOAD_FOLDER, exist_ok=True)

        file_path = os.path.join(Config.USER_SUMMARY_PPTX_DOWNLOAD_FOLDER, file_name)
        root.save(file_path)

        return file_path

    def create_itemized_summary_excel(self, request_body):
        """
        The function `create_itemized_summary_excel` creates an Excel file with an itemized summary
        based on the provided request body.

        Args:
          request_body: The `request_body` parameter is a list of dictionaries. Each dictionary
        represents a file record and contains the following keys:

        Returns:
          the file path of the created Excel file.
        """
        workbook = Workbook()
        sheet = workbook.active
        headers = ["Sl No.", "Title", "Summary"]
        sheet.append(headers)
        # adding data to workbook
        for file_record in request_body:
            row = [
                file_record.get("seq_number", 0) + 1,
                file_record.get("title", ""),
                file_record.get("summary", ""),
            ]
            sheet.append(row)

        timestamp = datetime.datetime.utcnow().timestamp()
        file_name = f"Itemized_summary_{timestamp}.xlsx"
        # Ensure that the user image upload folder exists
        os.makedirs(Config.USER_SUMMARY_XLSX_DOWNLOAD_FOLDER, exist_ok=True)

        file_path = os.path.join(Config.USER_SUMMARY_XLSX_DOWNLOAD_FOLDER, file_name)
        workbook.save(file_path)
        workbook.close()
        return file_path

    def create_highlights_summary_excel(self, request_body):
        """
        The function `create_highlights_summary_excel` creates an Excel file with a summary of
        highlights based on the provided request body.

        Args:
          request_body: The `request_body` parameter is a list of dictionaries. Each dictionary
        represents a file record and contains the following keys:

        Returns:
          the file path of the created Excel file.
        """
        print(request_body)
        workbook = Workbook()
        sheet = workbook.active
        # adding headers
        headers = ["Sl No.", "Title", "Sentences and Key phrases"]
        sheet.append(headers)
        sheet.cell(2, 3).value = "Sentence"
        sheet.cell(2, 4).value = "Key phrase"

        # adding data to workbook
        j = 3
        for i, file_record in enumerate(request_body):
            print(file_record)
            row = [file_record.get("seq_number", 0) + 1, file_record.get("title", "")]
            sheet.append(row)
            highlights = file_record["highlights"]
            for i, highlight in enumerate(highlights):
                print(highlight)
                sheet.cell(row=j, column=3).value = highlight.get("Sentence", "")
                sheet.cell(row=j, column=4).value = highlight.get("KeyPhrases", "")
                j += 1

            # print(row)

        timestamp = datetime.datetime.utcnow().timestamp()
        file_name = f"Highlights_summary_{timestamp}.xlsx"
        # Ensure that the user image upload folder exists
        os.makedirs(Config.USER_SUMMARY_XLSX_DOWNLOAD_FOLDER, exist_ok=True)

        file_path = os.path.join(Config.USER_SUMMARY_XLSX_DOWNLOAD_FOLDER, file_name)
        workbook.save(file_path)
        workbook.close()
        return file_path
        # workbook.save(f'app/excel_outputs/Output_{timestamp}.xlsx')
        # workbook.close()
        # return f"excel_outputs\Output_{timestamp}.xlsx"

    def delete_file(self, file_id, user_id):
        """
        The `delete_file` function deletes a file from a database and optionally from a cloud storage
        bucket, based on the file ID and user ID provided.

        Args:
          file_id: The `file_id` parameter is the unique identifier of the file that needs to be
        deleted. It is used to locate the file in the database.
          user_id: The `user_id` parameter is the unique identifier of the user who is requesting to
        delete the file.

        Returns:
          a boolean value indicating whether the file was successfully deleted or not.
        """
        m_db = MongoClient.connect()

        file = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].find_one(
            {"_id": ObjectId(file_id), "createdBy._id": ObjectId(user_id)}
        )

        root = file["root"]
        print(root)
        # our flag which is used to check if document has been deleted from gcp / system folder before deleting from the database
        deleted = False
        folder_substring = "/" + user_id + "/"
        if Config.GCP_PROD_ENV:
            file_name = file["virtualFileName"]
            if root == folder_substring:
                file_to_del = root[1:] + file_name
            else:
                file_to_del = root[1:] + "/" + file_name
            print("File to delete: " + file_to_del)
            bucket = Production.get_users_bucket()
            blob = bucket.blob(file_to_del)

            if blob.exists():
                blob.delete()
                deleted = True
                print(f"🟢 File '{file_name}' deleted successfully from bucket!.")
            else:
                deleted = False
                print(f"🔺 File '{file_name}' does NOT EXIST in bucket!.")
        else:
            file_path = MyDocumentsService.get_file_path(file, user_id)
            if os.path.isfile(file_path):
                # Delete the file
                os.remove(file_path)
                deleted = True
            else:
                deleted = False
                print("Invalid file path: %s" % file_path)

        # Delete document if the user is the one who uploaded the file
        if deleted == True:
            delete_response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].delete_one(
                {"_id": ObjectId(file_id), "createdBy._id": ObjectId(user_id)}
            )
            if delete_response.deleted_count > 0:
                return True

        # If the user has been shared the document then remove the user_id from the 'usersWithAccess' array
        print("Removing shared access!")
        user_access_update_response = m_db[
            Config.MONGO_DOCUMENT_MASTER_COLLECTION
        ].update_one(
            {"_id": ObjectId(file_id)},
            {"$pull": {"usersWithAccess": ObjectId(user_id)}},
        )

        return deleted

    def delete_files(self, file_ids, user_id):
        """
        The `delete_files` function deletes files from a database and also deletes the corresponding
        files from a storage system if the system is in production mode.

        Args:
          file_ids: The `file_ids` parameter is a list of file IDs that need to be deleted. Each file ID
        is a unique identifier for a file in the database.
          user_id: The `user_id` parameter is the unique identifier of the user who is requesting to
        delete the files.

        Returns:
          a boolean value indicating whether the files were successfully deleted or not.
        """
        m_db = MongoClient.connect()
        files_to_del = [ObjectId(file_id) for file_id in file_ids]
        # print("Files to delete : ", files_to_del)
        # print("User type : ", type(user_id))
        files = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].find(
            {"_id": {"$in": files_to_del}, "createdBy._id": ObjectId(user_id)}
        )
        # print(files)
        # for file in files:
        #     print("File : ", file)

        deleted = False

        for file in files:
            root = file["root"]
            print(root)
            root_folder_path = os.path.join(Config.USER_FOLDER, user_id)
            # our flag which is used to check if document has been deleted from gcp / system folder before deleting from the database

            folder_substring = "/" + user_id + "/"
            if Config.GCP_PROD_ENV:
                file_name = file["virtualFileName"]
                if root == folder_substring:
                    file_to_del = root[1:] + file_name
                else:
                    file_to_del = root[1:] + "/" + file_name
                print("File to delete: " + file_to_del)
                bucket = Production.get_users_bucket()
                blob = bucket.blob(file_to_del)

                if blob.exists():
                    blob.delete()
                    deleted = True
                    # print(
                    #     f"File '{file_name}' deleted successfully from bucket '{self.bucket_name}'."
                    # )
                else:
                    deleted = False
                    # print(
                    #     f"File '{file_name}' does not exist in bucket '{self.bucket_name}'."
                    # )
            else:
                file_name = file["virtualFileName"]
                # print("Filename : : : " + file_name)
                if root == folder_substring:
                    file_path = os.path.join(root_folder_path, file_name)
                    # Check if it's a file or a directory
                    if os.path.isfile(file_path):
                        print("File path : ", file_path)
                        # Delete the file
                        os.remove(file_path)
                        deleted = True
                    else:
                        print("Invalid file path: %s" % file_path)
                else:
                    new_root = root.replace(folder_substring, "")
                    print(new_root)
                    folder = os.path.join(root_folder_path, new_root)
                    file_path = os.path.join(folder, file_name)
                    # Check if it's a file or a directory
                    if os.path.isfile(file_path):
                        print("File path : ", file_path)
                        # Delete the file
                        os.remove(file_path)
                        deleted = True
                    else:
                        print("Invalid file path: %s" % file_path)

        # Delete document if the user is the one who uploaded the file
        if deleted == True:
            delete_response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].delete_many(
                {"_id": {"$in": files_to_del}, "createdBy._id": ObjectId(user_id)}
            )
            if delete_response.deleted_count > 0:
                return True

        # If the user has been shared the document then remove the user_id from the 'usersWithAccess' array
        print("Removing shared access!")
        user_access_update_response = m_db[
            Config.MONGO_DOCUMENT_MASTER_COLLECTION
        ].update_one(
            {"_id": {"$in": files_to_del}},
            {"$pull": {"usersWithAccess": ObjectId(user_id)}},
        )

        return bool(user_access_update_response.modified_count)

    def modify_document_shared_users(self, user_id, document_ids, target_user_ids=[]):
        """
        The function modifies the shared users of multiple documents by adding target user IDs to the
        usersWithAccess array.

        Args:
        user_id: The user_id parameter is the ID of the user who is modifying the documents.
        document_ids: The `document_ids` parameter is a list of unique identifiers of the documents that you
        want to modify the shared users for.
        target_user_ids: The `target_user_ids` parameter is a list of user IDs that you want to add to
        the shared users list of the documents.

        Returns:
        the number of documents modified in the database.
        """
        m_db = MongoClient.connect()
        modified_count = 0

        for document_id in document_ids:
            query = {"_id": ObjectId(document_id), "createdBy._id": ObjectId(user_id)}
            doc = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].find_one(query)

            if doc:
                update_query = {
                    "_id": ObjectId(document_id),
                    "createdBy._id": ObjectId(user_id),
                    "usersWithAccess": {"$nin": [ObjectId(target_user_id) for target_user_id in target_user_ids]},
                }

                update_statement = {
                    "$addToSet": {
                        "usersWithAccess": {
                            "$each": [ObjectId(target_user_id) for target_user_id in target_user_ids]
                        }
                    }
                }

                response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].update_one(update_query, update_statement)
                modified_count += response.modified_count

        return modified_count

    @staticmethod
    def _get_my_documents_pipeline():
        """Common My Documents pipeline

        Returns:
            list: Common My Documents pipeline
        """

        my_documents_pipeline = [
            PipelineStages.stage_add_fields(
                {
                    "createdBy._id": {"$toString": "$createdBy._id"},
                    "createdOn": {"$dateToString": {"date": "$createdOn"}},
                    "_id": {"$toString": "$_id"},
                    "virtualFileName": "$virtualFileName",
                    "usersWithAccess": {
                        "$map": {
                            "input": "$usersWithAccess",
                            "as": "user",
                            "in": {"$toString": "$$user"},
                        }
                    },
                }
            ),
            PipelineStages.stage_unset(["embeddings", "highlightsSummary"]),
        ]

        return my_documents_pipeline

    @staticmethod
    def _create_my_document_db_struct(title, description, filename, user, root):
        """
        The function `_create_my_document_db_struct` creates a document structure for a file in a
        document database.

        Args:
          title: The title of the document.
          description: The "description" parameter is a string that represents the description of the
        document. It provides additional information or details about the document.
          filename: The `filename` parameter is the name of the file that you want to create a document
        database structure for.
          user: The "user" parameter is a dictionary that represents the user who is creating the
        document. It contains information about the user, such as their ID and a reference to the user
        object.
          root: The "root" parameter is used to specify the root directory or folder where the document
        will be stored. If no root directory is provided (root == ""), the document will be stored in
        the root directory ("/").

        Returns:
          a document (doc) with various fields such as title, description, itemizedSummary,
        highlightsSummary, originalFileName, virtualFileName, createdBy, createdOn, embeddings, type,
        root, usersWithAccess, and storedOnCloud.
        """
        if root == "":
            root = "/"

        m_db = MongoClient.connect()
        # Check if the file with the same name already exists in the collection
        existing_files = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].find(
            {"root": root}
        )

        # Create a set of existing filenames for efficient lookup
        existing_filenames = set(file["originalFileName"] for file in existing_files)

        # If the filename already exists, increment the index count
        index = 1
        new_filename = filename
        while new_filename in existing_filenames:
            file_name_without_extension, file_extension = os.path.splitext(filename)
            print("File name without extension:", file_name_without_extension)
            new_filename = f"{file_name_without_extension}({index}){file_extension}"
            index = index + 1

        doc = {
            "title": title,
            "description": description,
            "itemizedSummary": "",  # update when itemized summary of this record is generated
            "highlightsSummary": "",  # update when highlight summary of this record is generated
            "originalFileName": new_filename,  # Use the unique filename
            "virtualFileName": "",
            "createdBy": {"_id": ObjectId(user["_id"]), "ref": "user"},
            "createdOn": datetime.datetime.utcnow(),
            "embeddings": None,
            "type": "File",
            "root": root,
            "usersWithAccess": [],
        }

        return doc

    def _my_documents_summarize_sumy_text_rank(
        self, document_ids: list, sentence_count: int
    ):
        """
        Summarizes the descriptions of user-uploaded documents one-by-one using Text Rank method from sumy module
        Sends them back one by one
        """
        m_db = MongoClient.connect()

        # Convert kiIds to a list of ObjectIDs
        document_ids = list(map(ObjectId, document_ids))

        documents = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].find(
            {"_id": {"$in": document_ids}}
        )

        summarizer = TextRankSummarizer()

        for idx, document in enumerate(documents):
            document_description = ""
            document_title = document["title"]
            document_created_date = str(document["createdOn"])

            if not document:
                yield document_description, idx, document_title, document_created_date, False

            document_description = document["description"]

            # If the summary of the document was already generated and stored then return that summary
            itemized_summary = document["itemizedSummary"]
            if len(itemized_summary) > 0:
                print("Found summary from DB!")
                text_summary = itemized_summary

            else:
                print("Generating summary!")

                # Otherwise generate the summary and store it
                parser = PlaintextParser.from_string(
                    document_description, Tokenizer("english")
                )

                summary = summarizer(parser.document, sentence_count)
                text_summary = ""
                for sentence in summary:
                    text_summary += str(sentence)

                # Store the text_summary for next time quick access
                m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].update_one(
                    {"_id": ObjectId(document["_id"])},
                    {"$set": {"itemizedSummary": text_summary}},
                )

            # yield text_summary, True
            yield text_summary, idx, document_title, document_created_date, True

    def _parse_pdf(self, file, original_filename, user, root):
        """
        The _parse_pdf function extracts content from the file and inserts a new record corresponding to the file.
            Args:
                file (file storage) : The file to be parsed
                original_filename (str) : The name of this file to be parsed
                user (dict): Identifies the user uploading the file.

        Args:
            self: Represent the instance of the class
            file : The file to be parsed
            original_filename (str) : The name of this file to be parsed
            user (str): Corresponds to the user uploading the file.

        Returns:
            The Objectid of the newly inserted record

        """

        m_db = MongoClient.connect()

        reader = PyPDF2.PdfReader(file)
        # This returns the document title
        # print(reader.metadata.get('/Title'))
        # first_page = reader.pages[0]
        # title = first_page.extract_text().strip()

        all_pages = reader.pages[0:]
        data = ""
        for page in all_pages:
            data += page.extract_text()
        # print("TITLE : ", title , "\nDATA : " , data)

        paragraphs = data.split("\n")  # Split content into paragraphs
        first_paragraph = paragraphs[0].strip() if paragraphs else ""
        # Assuming the first paragraph is the title
        if first_paragraph == "":
            first_paragraph = reader.metadata.get("/Title")

        file_data = MyDocumentsService._create_my_document_db_struct(
            first_paragraph, data, original_filename, user, root
        )
        response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].insert_one(file_data)

        if response:
            return str(response.inserted_id)

        return None

    def _parse_doc(self, file, filename, user, root):
        """
        The _parse_doc function extracts content from the file and inserts a new record corresponding to the file.
            Args:
                file (file storage) : The file to be parsed
                original_filename (str) : The name of this file to be parsed
                user (dict): Identifies the user uploading the file.

        Args:
            self: Represent the instance of the class
            file : The file to be parsed
            original_filename (str) : The name of this file to be parsed
            user (str): Corresponds to the user uploading the file.

        Returns:
            The ObjectId of the newly inserted record
            :param user:
            :param file:
            :param filename:

        """

        try:
            # print("File : ", file)
            document = Document(file)
            print("document : ", document)
            title = document.paragraphs[0].text.strip()
            if document.paragraphs[0] is not None:
                title = document.paragraphs[0].text.strip()
            else:
                title = ""
            data = ""
            for paragraph in document.paragraphs[1:]:
                data += paragraph.text

            file_data = MyDocumentsService._create_my_document_db_struct(
                title, data, filename, user, root
            )

            # print("File data  : ", file_data)

            m_db = MongoClient.connect()
            response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].insert_one(
                file_data
            )

            if response:
                return str(response.inserted_id)
            else:
                return None
        except Exception as e:
            print("Error uploading file : " + str(e))
            return None

    def _parse_pptx(self, file, filename, user, root):
        """
        The _parse_pptx function extracts content from the file and inserts a new record corresponding to the file.
            Args:
                file (file storage) : The file to be parsed
                original_filename (str) : The name of this file to be parsed
                user (dict): Identifies the user uploading the file.

        Args:
            self: Represent the instance of the class
            file : The file to be parsed
            original_filename (str) : The name of this file to be parsed
            user (str): Corresponds to the user uploading the file.

        Returns:
            The Objectid of the newly inserted record
            :param user:
            :param file:
            :param filename:

        """

        m_db = MongoClient.connect()

        ppt = Presentation(file)
        slides = ppt.slides
        first_slide = slides[0]
        if first_slide.shapes.title is not None:
            title = first_slide.shapes.title.text.strip()
        else:
            title = "No Title"

        data = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        data += paragraph.text

        file_data = MyDocumentsService._create_my_document_db_struct(
            title, data, filename, user, root
        )

        response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].insert_one(file_data)

        if response:
            return str(response.inserted_id)

        return None

    def _parse_text(self, file, filename, user, root):
        """
        The _parse_doc function extracts content from the file and inserts a new record corresponding to the file.
            Args:
                file (file storage) : The file to be parsed
                original_filename (str) : The name of this file to be parsed
                user (dict): Identifies the user uploading the file.

        Args:
            self: Represent the instance of the class
            file : The file to be parsed
            original_filename (str) : The name of this file to be parsed
            user (str): Corresponds to the user uploading the file.

        Returns:
            The ObjectId of the newly inserted record
            :param user:
            :param file:
            :param filename:

        """

        m_db = MongoClient.connect()
        print(file)

        try:
            title = file.readline().decode(
                "utf-8"
            )  # Decode the bytes into a string using the appropriate encoding
            content = file.read().decode("utf-8")
            name, extension = os.path.splitext(filename)
            # print(f"name {name}, ext {extension}, file {filename}")
            if (
                title == "" or title == "\r\n"
            ):  # if there is a space or new line here \r\n will match the new line
                title = name

            file_data = MyDocumentsService._create_my_document_db_struct(
                title, content, filename, user, root
            )
            print(file_data)
            response = m_db[Config.MONGO_DOCUMENT_MASTER_COLLECTION].insert_one(
                file_data
            )

            if response:
                return str(response.inserted_id)

        except Exception as e:
            print("An error occurred:", e)
            return None

    def share_document_via_email(self, user_id, document_ids, email_ids, subject: str, message: str):
        try:
            files = self.get_files_by_id(document_ids)
            virtual_file_names = [file["virtualFileName"] for file in files]
            user = UserService.get_user_by_id(user_id)
            file_details = [MyDocumentsService().get_file_contents(virtual_file_name) for virtual_file_name in virtual_file_names]
            file_contents, file_names = zip(*file_details) 
            attachments = [{"content": get_base64_encoding(file_content), "name": file_name} for (file_content, file_name) in zip(file_contents, file_names)]
            recipients = [{"name": None, "email": email_id} for email_id in email_ids]
                        
            email_response = send_mail(
                subject = subject or Config.DEFAULT_DOCUMENT_EMAIL_SUBJECT,
                htmlMailBody = message or Config.DEFAULT_DOCUMENT_EMAIL_MESSAGE,
                recipients = recipients,
                sender = {"name": user["name"], "email": user["email"]},
                attachments = attachments
            )
            
            return email_response
                
        except Exception as e:
            Common.exception_details("MyDocumentsService.share_document_via_email", e)
            return False
    
    @staticmethod
    def get_file_contents(virtual_document_name):
        try:
            file = MyDocumentsService().get_file_by_virtual_name(virtual_document_name)
            file_created_by = str(file["createdBy"]["_id"])
            file_root = str(file["root"])

            if Config.GCP_PROD_ENV:
                bucket = Production.get_users_bucket()
                path = file_root[1:] + ("/" if file_root != f"/{file_created_by}/" else "")
                blob = bucket.blob(path + virtual_document_name)
                bytes = blob.download_as_bytes()
                return io.BytesIO(bytes), file["originalFileName"]
            else:
                user_folder_path = os.path.join(Config.USER_FOLDER, file_root[1:])
                file_save_path = os.path.join(user_folder_path, virtual_document_name)
                with open(file_save_path, 'rb') as file_handle:
                    file_bytes = file_handle.read()
                    return io.BytesIO(file_bytes), file["originalFileName"]
                    
        except Exception as e:
            Common.exception_details("MyDocumentsService.get_file_contents", e)
            return None, None
            
        